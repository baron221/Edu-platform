from pptx import Presentation
from pptx.util import Inches
import os

def create_premium_pptx():
    prs = Presentation()
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_images = [
        "slide1_1778495574249.png",
        "slide2_1778495588209.png",
        "slide3_1778495599531.png",
        "slide4_1778495607748.png",
        "slide5_1778495615777.png",
        "slide6_1778495623525.png",
        "slide7_1778495631296.png"
    ]

    base_path = r"C:\Users\baron\.gemini\antigravity\brain\9cde4dd5-c16b-4eef-8c92-aeb0c932585f"

    for img_name in slides_images:
        img_path = os.path.join(base_path, img_name)
        if os.path.exists(img_path):
            # Use blank layout
            slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(slide_layout)
            
            # Add picture to fill the whole slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            print(f"Warning: {img_path} not found.")

    save_path = r"c:\Users\baron\Documents\Edu-platform\Edu-platform\edunation\EduNation_Premium_Presentation.pptx"
    prs.save(save_path)
    print(f"Premium PPTX successfully generated at {save_path}")

if __name__ == "__main__":
    create_premium_pptx()
